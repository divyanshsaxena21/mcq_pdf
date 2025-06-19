import re
import textwrap
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi

def clean_text(text):
    lines = text.split("\n")
    return "\n".join([line for line in lines if not re.match(r"^\s*(page \d+|[\d\s\-]+)$", line.strip().lower())])

def chunk_text(text, max_chars=1000, min_len=100):
    cleaned = re.sub(r"\s*\n\s*", " ", text)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    chunks = textwrap.wrap(cleaned, width=max_chars, break_long_words=False, break_on_hyphens=False)
    return [chunk.strip() for chunk in chunks if len(chunk.strip()) > min_len]

def extract_sections_from_pdf(path):
    doc = fitz.open(path)
    text = "\n".join([clean_text(page.get_text()) for page in doc])
    return chunk_text(text)

def extract_sections_from_docx(path):
    doc = Document(path)
    text = " ".join([para.text for para in doc.paragraphs])
    return chunk_text(text)

def extract_sections_from_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return chunk_text(text)

def extract_video_id(url):
    import re
    patterns = [r"youtu\.be/([^?&]+)", r"youtube\.com/watch\?v=([^&]+)", r"youtube\.com/embed/([^?&]+)"]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            return match.group(1)
    return None

def extract_sections_from_youtube(url):
    video_id = extract_video_id(url)
    if not video_id:
        raise ValueError("Invalid YouTube URL")
    transcript = YouTubeTranscriptApi.get_transcript(video_id)
    text = " ".join([entry['text'] for entry in transcript])
    return chunk_text(text)
